"""Tests for PowerPoint -> Markdown conversion."""

from pathlib import Path

import pytest


def _make_sample_pptx(path: Path) -> None:
    """Create a minimal one-slide PPTX with known text using python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    # 제목 추가
    if hasattr(slide.shapes, "title") and slide.shapes.title:
        slide.shapes.title.text = "Test Slide Title"

    # 텍스트 추가
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = "Body text for thomas_utils PPTX conversion."

    prs.save(str(path))


def test_convert_pptx_non_empty(tmp_path: Path) -> None:
    """convert_pptx() returns non-empty markdown for a minimal PPTX."""
    from thomas_utils.converters import convert_pptx

    pptx_path = tmp_path / "sample.pptx"
    _make_sample_pptx(pptx_path)
    result = convert_pptx(str(pptx_path))
    assert result
    assert "**Type**" in result
    assert "### Content" in result
    assert "**Title**" in result or "Test Slide Title" in result
    assert "Test" in result or "Slide" in result or "Body" in result or "thomas_utils" in result


def test_convert_pptx_missing_file() -> None:
    """convert_pptx() raises FileNotFoundError for missing path."""
    from thomas_utils.converters import convert_pptx

    with pytest.raises(FileNotFoundError, match="not found"):
        convert_pptx("/nonexistent/sample.pptx")


def test_convert_pptx_invalid_extension(tmp_path: Path) -> None:
    """convert_pptx() raises ValueError for non-pptx file."""
    from thomas_utils.converters import convert_pptx

    pdf_path = tmp_path / "sample.pdf"
    pdf_path.write_bytes(b"fake pdf")
    with pytest.raises(ValueError, match="Expected .pptx"):
        convert_pptx(str(pdf_path))
