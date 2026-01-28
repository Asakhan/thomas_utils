"""Tests for CLI."""

from pathlib import Path

import pytest


def _make_sample_pdf(path: Path) -> None:
    import pymupdf

    doc = pymupdf.open()
    page = doc.new_page()
    page.insert_text((72, 72), "CLI test content.")
    doc.save(str(path))
    doc.close()


def test_cli_pdf2md(tmp_path: Path) -> None:
    """pdf2md writes UTF-8 Markdown and exits 0."""
    from thomas_utils.cli import _pdf2md
    from thomas_utils.converters import convert

    pdf_path = tmp_path / "in.pdf"
    _make_sample_pdf(pdf_path)
    out_path = tmp_path / "out.md"

    class Args:
        input = str(pdf_path)
        output = str(out_path)
        pages = None
        engine = "pymupdf"

    code = _pdf2md(Args())
    assert code == 0
    assert out_path.exists()
    text = out_path.read_text(encoding="utf-8")
    assert text
    assert "CLI" in text or "test" in text or "content" in text


def test_cli_pdf2md_default_output_dir(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    """When output is not specified, result is written to output/INPUT_NAME.md under cwd."""
    from thomas_utils.cli import _pdf2md

    monkeypatch.chdir(tmp_path)
    pdf_path = tmp_path / "doc.pdf"
    _make_sample_pdf(pdf_path)

    class Args:
        input = str(pdf_path)
        output = None
        pages = None
        engine = "pymupdf"

    code = _pdf2md(Args())
    assert code == 0
    out_path = tmp_path / "output" / "doc.md"
    assert out_path.exists()
    text = out_path.read_text(encoding="utf-8")
    assert "CLI" in text or "test" in text or "content" in text


def _make_sample_pptx(path: Path) -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if hasattr(slide.shapes, "title") and slide.shapes.title:
        slide.shapes.title.text = "CLI test slide"
    prs.save(str(path))


def test_cli_pptx2md(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    """pptx2md writes UTF-8 Markdown and exits 0."""
    from thomas_utils.cli import _pptx2md

    monkeypatch.chdir(tmp_path)
    pptx_path = tmp_path / "in.pptx"
    _make_sample_pptx(pptx_path)
    out_path = tmp_path / "out.md"

    class Args:
        input = str(pptx_path)
        output = str(out_path)
        slides = None

    code = _pptx2md(Args())
    assert code == 0
    assert out_path.exists()
    text = out_path.read_text(encoding="utf-8")
    assert text
    assert "CLI" in text or "test" in text or "slide" in text


def test_cli_pptx2md_default_output_dir(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    """When output is not specified, pptx2md writes to output/INPUT_NAME.md under cwd."""
    from thomas_utils.cli import _pptx2md

    monkeypatch.chdir(tmp_path)
    pptx_path = tmp_path / "presentation.pptx"
    _make_sample_pptx(pptx_path)

    class Args:
        input = str(pptx_path)
        output = None
        slides = None

    code = _pptx2md(Args())
    assert code == 0
    out_path = tmp_path / "output" / "presentation.md"
    assert out_path.exists()
    text = out_path.read_text(encoding="utf-8")
    assert "CLI" in text or "test" in text or "slide" in text
