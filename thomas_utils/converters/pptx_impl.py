"""PowerPoint (.pptx) -> Markdown conversion using python-pptx."""

import base64
import json
import os
import re
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import List, Optional, Union

# #region agent log
LOG_PATH = Path(__file__).resolve().parents[2] / ".cursor" / "debug.log"
SESSION = "debug-session"
RUN_ID = "pptx-import-check"
try:
    import importlib.util
    spec = importlib.util.find_spec("pptx")
    pptx_available = spec is not None
    payload = {
        "sessionId": SESSION,
        "runId": RUN_ID,
        "hypothesisId": "H1",
        "location": "pptx_impl.py:import",
        "message": "pptx module availability check",
        "data": {"pptx_available": pptx_available, "sys_executable": sys.executable},
        "timestamp": __import__("time").time() * 1000,
    }
    LOG_PATH.parent.mkdir(parents=True, exist_ok=True)
    with open(LOG_PATH, "a", encoding="utf-8") as f:
        f.write(json.dumps(payload) + "\n")
except Exception:
    pass
# #endregion

# 마크다운 이미지 문법 줄 제거용 (슬라이드 텍스트에 포함된 경우 제외)
_IMAGE_LINE_PATTERN = re.compile(r"^!\[.*\]\(.*\)\s*$", re.MULTILINE)

# 코드블록 후보: 문단이 코드처럼 보이는 패턴
_CODE_LINE_PATTERN = re.compile(
    r"^(?:from\s+\S+\s+import|import\s+\S+|def\s+\w+|class\s+\w+|\s{2,}\S)"
)


def _strip_image_lines(text: str) -> str:
    """Remove lines that are markdown image references (e.g. ![...](...))."""
    return _IMAGE_LINE_PATTERN.sub("", text).strip()


try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
except ImportError as e:
    raise ImportError(
        "python-pptx is not installed. Please run: pip install python-pptx"
    ) from e


def _content_shape_sort_key(shape) -> tuple:
    """Sort key for visual order: top then left (reading order)."""
    top = getattr(shape, "top", 0) or 0
    left = getattr(shape, "left", 0) or 0
    return (top, left)


def _slide_type_from_layout_name(name: Optional[str]) -> str:
    """Map slide layout name to Type string."""
    if not name:
        return "Content Slide"
    n = (name or "").strip().lower()
    if "title" in n and ("content" not in n and "body" not in n and "object" not in n):
        if "section" in n or "header" in n or "divider" in n or "구역" in n:
            return "Section Divider"
        return "Title Slide"
    if "section" in n or "header" in n or "구역" in n:
        return "Section Divider"
    return "Content Slide"


def _layout_hint_from_layout_name(name: Optional[str]) -> Optional[str]:
    """Layout description from layout name (e.g. Center-aligned)."""
    if not name or "center" not in (name or "").lower():
        return None
    return "Center-aligned"


def _get_placeholder_type(shape) -> Optional[int]:
    """Return PP_PLACEHOLDER type for shape if it is a placeholder."""
    try:
        if getattr(shape, "is_placeholder", False) and getattr(shape, "placeholder_format", None):
            return getattr(shape.placeholder_format, "type", None)
    except Exception:
        pass
    return None


def _is_content_shape(shape, title: Optional[str], subtitle: Optional[str]) -> bool:
    """True if shape contributes to Content (body, table, text, or picture slot for ordering)."""
    pph = _get_placeholder_type(shape)
    if pph is not None:
        if pph in (
            PP_PLACEHOLDER.TITLE,
            PP_PLACEHOLDER.CENTER_TITLE,
            getattr(PP_PLACEHOLDER, "VERTICAL_TITLE", None),
            PP_PLACEHOLDER.SUBTITLE,
        ):
            return False
        if pph in (PP_PLACEHOLDER.BODY, getattr(PP_PLACEHOLDER, "VERTICAL_BODY", None)):
            return True
        return False
    if getattr(shape, "has_table", False):
        return True
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
        return True
    if hasattr(shape, "text_frame") and shape.text_frame:
        return True
    if hasattr(shape, "text") and (shape.text or "").strip():
        return True
    return False


def _table_to_markdown(table) -> str:
    """Convert python-pptx table to markdown table string. Uses Pandas if available, else fallback."""
    rows: List[List[str]] = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cells.append((cell.text or "").replace("|", "\\|").replace("\n", " ").strip())
        rows.append(cells)
    if not rows:
        return ""
    col_count = max(len(r) for r in rows)
    for r in rows:
        while len(r) < col_count:
            r.append("")
    try:
        import pandas as pd
        df = pd.DataFrame(rows)
        return df.to_markdown(index=False, tablefmt="github")
    except ImportError:
        pass
    lines = ["| " + " | ".join(rows[0]) + " |", "| " + " | ".join("---" for _ in range(col_count)) + " |"]
    for r in rows[1:]:
        lines.append("| " + " | ".join(r) + " |")
    return "\n".join(lines)


def _extract_omml_from_shape(shape) -> List[str]:
    """Extract OMML (Office Math) XML from shape for LaTeX conversion. Returns list of LaTeX strings (empty if no math or converter missing)."""
    try:
        el = getattr(shape, "_element", None)
        if el is None:
            return []
        # python-pptx uses lxml; namespace for Office Math
        ns = "http://schemas.openxmlformats.org/officeDocument/2006/math"
        omaths = el.findall(f".//{{{ns}}}oMath") if hasattr(el, "findall") else []
        if not omaths:
            return []
        out: List[str] = []
        for omath in omaths:
            latex = _omml_to_latex(omath)
            if latex:
                out.append(latex)
        return out
    except Exception:
        return []


def _omml_to_latex(omml_element) -> str:
    """Convert one OMML element to LaTeX. Requires optional [pptx-math] extra; otherwise returns empty string."""
    try:
        from officemath2latex import convert as omml_convert
        return (omml_convert(omml_element) or "").strip()
    except ImportError:
        pass
    except Exception:
        pass
    return ""


def _text_frame_to_structured_content(text_frame) -> str:
    """Convert text_frame paragraphs to markdown (lists by level, code blocks, plain)."""
    parts: List[str] = []
    current_code_lines: List[str] = []
    code_indicators = ("from ", "import ", "def ", "class ")

    def flush_code():
        nonlocal current_code_lines
        if current_code_lines:
            block = "\n".join(current_code_lines)
            parts.append("```\n" + block + "\n```")
            current_code_lines = []

    for para in text_frame.paragraphs:
        text = (para.text or "").strip()
        if not text:
            flush_code()
            continue
        level = getattr(para, "level", 0) or 0
        # 코드 유사: 앞쪽 공백 2칸 이상 또는 from/import/def/class 로 시작
        looks_like_code = (
            text.startswith("  ") or
            any(text.startswith(k) for k in code_indicators) or
            _CODE_LINE_PATTERN.match(text) is not None
        )
        if looks_like_code and level == 0:
            current_code_lines.append(text)
            continue
        flush_code()
        if level > 0:
            indent = "   " * level
            parts.append(indent + "- " + text)
        else:
            parts.append(text)
    flush_code()
    return "\n\n".join(p for p in parts if p.strip())


def _structure_body_content(
    body_text_parts: List[str],
    table_md_parts: List[str],
    other_text_parts: List[str],
) -> str:
    """Combine body text, table markdown, and other text into one Content block."""
    segments: List[str] = []
    if body_text_parts:
        combined = "\n\n".join(body_text_parts)
        combined = _strip_image_lines(combined)
        if combined.strip():
            segments.append(combined)
    for t in table_md_parts:
        if t.strip():
            segments.append(t)
    for t in other_text_parts:
        if t.strip():
            segments.append(_strip_image_lines(t))
    return "\n\n".join(segments).strip()


def convert(
    pptx_path: Union[str, Path],
    slides: Optional[List[int]] = None,
    use_llm: bool = False,
    engine: str = "python-pptx",
    use_llm_multimodal: bool = False,
) -> str:
    """Convert PowerPoint to structured Markdown.

    Each slide is emitted with Type, Layout, Title, Subtitle, and Content
    (tables, lists, code blocks, and plain paragraphs).

    Args:
        pptx_path: Path to the PPTX file.
        slides: Ignored (always converts all slides for now).
        use_llm: If True, run optional LLM polish on the result (requires pptx-llm extra).
        engine: "python-pptx" (default) or "unstructured" (requires [unstructured] extra).
        use_llm_multimodal: If True, render each slide to image and convert via vision LLM (GPT-4o).

    Returns:
        UTF-8 Markdown string.
    """
    if use_llm_multimodal:
        return _convert_pptx_multimodal(pptx_path, use_llm=use_llm)
    if engine == "unstructured":
        from thomas_utils.converters.pptx_unstructured_impl import convert_unstructured
        result = convert_unstructured(pptx_path)
        if use_llm:
            result = _llm_polish(result)
        return result
    path = Path(pptx_path)
    if not path.exists():
        raise FileNotFoundError(f"PPTX not found: {path}")
    if not path.suffix.lower() == ".pptx":
        raise ValueError(f"Expected .pptx file, got: {path}")

    prs = Presentation(str(path))
    md_parts: List[str] = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_layout = getattr(slide, "slide_layout", None)
        layout_name = getattr(slide_layout, "name", None) if slide_layout else None
        slide_type = _slide_type_from_layout_name(layout_name)
        layout_hint = _layout_hint_from_layout_name(layout_name)

        title = None
        subtitle = None

        # 1) Title/Subtitle from placeholders only
        for shape in slide.shapes:
            pph = _get_placeholder_type(shape)
            if pph is None:
                continue
            if pph in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE, getattr(PP_PLACEHOLDER, "VERTICAL_TITLE", None)):
                if hasattr(shape, "text") and shape.text.strip():
                    title = shape.text.strip()
            elif pph == PP_PLACEHOLDER.SUBTITLE:
                if hasattr(shape, "text") and shape.text.strip():
                    subtitle = shape.text.strip()

        # 2) Content shapes in visual order (Top, then Left)
        content_shapes = [s for s in slide.shapes if _is_content_shape(s, title, subtitle)]
        content_shapes.sort(key=_content_shape_sort_key)
        content_segments: List[str] = []

        for shape in content_shapes:
            # Shape decomposition: table, picture, text_frame (수식은 별도 단계에서 처리)
            if getattr(shape, "has_table", False) and shape.table:
                content_segments.append(_table_to_markdown(shape.table))
                continue
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                # 이미지 미포함 정책; --include-images 시 여기서 분기 가능
                continue
            if hasattr(shape, "text_frame") and shape.text_frame:
                text = (shape.text_frame.text or "").strip()
                if not text:
                    continue
                if title and text == title or subtitle and text == subtitle:
                    continue
                structured = _text_frame_to_structured_content(shape.text_frame)
                omml_latex = _extract_omml_from_shape(shape)
                if omml_latex:
                    structured = (structured or "") + "\n\n" + "\n\n".join(f"$${l}$$" for l in omml_latex)
                if structured:
                    content_segments.append(structured)
                continue
            if hasattr(shape, "text") and shape.text.strip():
                text = _strip_image_lines(shape.text.strip())
                if text and text != title and text != subtitle:
                    content_segments.append(text)

        content_block = "\n\n".join(
            _strip_image_lines(s).strip() for s in content_segments if s.strip()
        ).strip()

        # Build slide block per plan template
        block_lines = [
            f"## Slide {slide_idx + 1}",
            f"**Type**: {slide_type}",
        ]
        if layout_hint:
            block_lines.append(f"**Layout**: {layout_hint}")
        if title:
            block_lines.append(f"**Title**: {title}")
        if subtitle:
            block_lines.append(f"**Subtitle**: {subtitle}")
        block_lines.append("")
        if content_block:
            block_lines.append("### Content")
            block_lines.append("")
            block_lines.append(content_block)
        else:
            block_lines.append("### Content")
            block_lines.append("")

        md_parts.append("\n".join(block_lines))
        if slide_idx < len(prs.slides) - 1:
            md_parts.append("\n---\n\n")

    result = "\n".join(md_parts)
    result = _IMAGE_LINE_PATTERN.sub("", result)
    result = re.sub(r"\n{3,}", "\n\n", result).strip()
    result = result + "\n" if result else result

    if use_llm:
        result = _llm_polish(result)
    return result


def _render_pptx_slides_to_images(pptx_path: Union[str, Path]) -> List[bytes]:
    """Render each PPTX slide to PNG image bytes. Tries Windows PowerPoint COM, then LibreOffice + PyMuPDF."""
    path = Path(pptx_path).resolve()
    if not path.exists():
        raise FileNotFoundError(f"PPTX not found: {path}")
    if path.suffix.lower() != ".pptx":
        raise ValueError(f"Expected .pptx file, got: {path}")

    # 1) Windows: PowerPoint COM (pywin32)
    # Do not set app.Visible = 0: some Office configs raise "Hiding the application window is not allowed."
    try:
        import win32com.client
        app = win32com.client.Dispatch("PowerPoint.Application")
        path_str = os.path.normpath(str(path))
        prs = app.Presentations.Open(path_str, WithWindow=False)
        n = prs.Slides.Count
        out: List[bytes] = []
        with tempfile.TemporaryDirectory() as tmp:
            for i in range(1, n + 1):
                png_path = Path(tmp) / f"slide_{i}.png"
                prs.Slides(i).Export(str(png_path), "PNG")
                out.append(png_path.read_bytes())
        prs.Close()
        app.Quit()
        return out
    except ModuleNotFoundError as _e:
        raise RuntimeError(
            "멀티모달(슬라이드 이미지)을 쓰려면 Windows에서 pywin32가 필요합니다: pip install pywin32. "
            "PowerPoint가 설치되어 있어야 합니다."
        ) from _e
    except Exception as _e:
        sys.stderr.write(
            f"PowerPoint COM 실패 ({type(_e).__name__}: {_e}), LibreOffice 경로로 시도합니다.\n"
        )
        pass

    # 2) Fallback: LibreOffice -> PDF, then PyMuPDF -> PNG per page
    try:
        import fitz
    except ImportError:
        fitz = None
    if fitz is None:
        raise RuntimeError(
            "Slide-to-image requires either: (1) Windows + PowerPoint + pywin32, or "
            "(2) LibreOffice (soffice) in PATH + PyMuPDF. Install: pip install pywin32 (Windows) or pymupdf, and ensure LibreOffice is installed."
        )
    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        try:
            result = subprocess.run(
                [
                    "soffice",
                    "--headless",
                    "--convert-to", "pdf",
                    "--outdir", str(tmp_path),
                    str(path),
                ],
                capture_output=True,
                timeout=120,
            )
        except FileNotFoundError as _e2:
            raise FileNotFoundError(
                "soffice (LibreOffice)를 PATH에서 찾을 수 없습니다. "
                "Windows에서는 PowerPoint를 쓰려면: pip install pywin32 를 설치하고 PowerPoint가 설치되어 있어야 합니다. "
                "또는 LibreOffice를 설치한 뒤 'soffice'가 PATH에 있도록 하세요."
            ) from _e2
        except Exception as _e2:
            raise
        if result.returncode != 0:
            raise RuntimeError(
                "LibreOffice conversion failed. Install LibreOffice and ensure 'soffice' is in PATH."
            )
        pdf_path = tmp_path / (path.stem + ".pdf")
        if not pdf_path.exists():
            raise RuntimeError("LibreOffice did not produce PDF.")
        doc = fitz.open(pdf_path)
        out = []
        for page in doc:
            pix = page.get_pixmap(alpha=False)
            out.append(pix.tobytes("png"))
        doc.close()
        return out


def _llm_slide_image_to_md(image_bytes: bytes, slide_index: int) -> str:
    """Convert a single slide image to markdown via multimodal LLM (GPT-4o). Uses OPENAI_API_KEY from .env."""
    try:
        import os
        from dotenv import load_dotenv
        load_dotenv()
        api_key = os.environ.get("OPENAI_API_KEY")
        if not api_key:
            return f"## Slide {slide_index + 1}\n**Type**: Content Slide\n\n### Content\n\n<!-- OPENAI_API_KEY not set -->\n\n"
        import openai
    except ImportError:
        return f"## Slide {slide_index + 1}\n**Type**: Content Slide\n\n### Content\n\n<!-- openai or python-dotenv not installed -->\n\n"
    b64 = base64.b64encode(image_bytes).decode("ascii")
    prompt = (
        "이 슬라이드 이미지를 마크다운으로 변환해줘. 다음 형식만 사용하고 마크다운만 출력해.\n\n"
        f"## Slide {slide_index + 1}\n"
        "**Type**: (Title Slide | Content Slide | Section Divider 중 하나)\n"
        "**Title**: (제목이 있으면)\n"
        "**Subtitle**: (부제가 있으면)\n"
        "### Content\n"
        "(본문: 표는 마크다운 테이블, 리스트는 -, 코드는 ``` 블록으로)"
    )
    client = getattr(openai, "OpenAI", None)
    if not client:
        return f"## Slide {slide_index + 1}\n**Type**: Content Slide\n\n### Content\n\n\n\n"
    try:
        c = client(api_key=api_key)
        r = c.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}"}},
                    ],
                }
            ],
            max_tokens=4096,
        )
        if r.choices and r.choices[0].message.content:
            return r.choices[0].message.content.strip() + "\n"
    except Exception:
        pass
    return f"## Slide {slide_index + 1}\n**Type**: Content Slide\n\n### Content\n\n\n\n"


def _convert_pptx_multimodal(
    pptx_path: Union[str, Path],
    use_llm: bool = False,
) -> str:
    """Convert PPTX to Markdown by rendering each slide to image and calling vision LLM (GPT-4o)."""
    images = _render_pptx_slides_to_images(pptx_path)
    md_parts: List[str] = []
    for i, img_bytes in enumerate(images):
        md_parts.append(_llm_slide_image_to_md(img_bytes, i))
        if i < len(images) - 1:
            md_parts.append("\n---\n\n")
    result = "\n".join(md_parts)
    result = re.sub(r"\n{3,}", "\n\n", result).strip()
    result = result + "\n" if result else result
    if use_llm:
        result = _llm_polish(result)
    return result


def _llm_polish(md: str) -> str:
    """Optional LLM polish: naturalize wording, add code block language, etc. Uses OPENAI_API_KEY from .env."""
    try:
        import os
        from dotenv import load_dotenv
        load_dotenv()
        api_key = os.environ.get("OPENAI_API_KEY")
        if not api_key:
            return md
        import openai
    except ImportError:
        return md
    client = getattr(openai, "OpenAI", None)
    if not client:
        return md
    try:
        c = client(api_key=api_key)
        r = c.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {
                    "role": "user",
                    "content": "아래는 PPT에서 추출한 마크다운이다. 슬라이드 재구성에 쓸 수 있도록, 형식(## Slide, Type, Layout, Title, Subtitle, Content)은 유지한 채로 문장만 자연스럽게 다듬고, 표 제목·코드블록 언어는 필요 시 보완해라. 마크다운만 출력해라.\n\n" + md,
                }
            ],
        )
        if r.choices and r.choices[0].message.content:
            return r.choices[0].message.content.strip() + "\n"
    except Exception:
        pass
    return md
